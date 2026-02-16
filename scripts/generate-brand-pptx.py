"""
Generate a branded PowerPoint template for Rent-A-Vacation.
Run: python scripts/generate-brand-pptx.py
Output: docs/RAV-Brand-Template.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
TEAL = RGBColor(0x1C, 0x72, 0x68)
CORAL = RGBColor(0xE8, 0x70, 0x3A)
CREAM = RGBColor(0xF8, 0xF6, 0xF3)
NAVY = RGBColor(0x1D, 0x2E, 0x38)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SAND = RGBColor(0xF0, 0xEB, 0xE3)
MUTED = RGBColor(0x6B, 0x7B, 0x85)
SUCCESS = RGBColor(0x1F, 0xA6, 0x6E)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT


def add_bg_rect(slide, color, left=0, top=0, width=None, height=None):
    """Add a background rectangle."""
    w = width or SLIDE_WIDTH
    h = height or SLIDE_HEIGHT
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=NAVY, bold=False, alignment=PP_ALIGN.LEFT, font_name="Roboto"):
    """Add a text box with specified formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_logo_text(slide, left, top, size=14):
    """Add 'RENT-A-VACATION' text as logo placeholder."""
    add_text_box(slide, left, top, Inches(3), Inches(0.5),
                 "RENT-A-VACATION", font_size=size, color=WHITE, bold=True)


def add_accent_bar(slide, top, width=Inches(2), color=CORAL):
    """Add a thin accent bar."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), top, width, Inches(0.06))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_footer(slide, dark_bg=False):
    """Add footer with brand name and tagline."""
    fg = WHITE if dark_bg else MUTED
    add_text_box(slide, Inches(0.8), Inches(6.8), Inches(6), Inches(0.4),
                 "rent-a-vacation.com  |  Name Your Price. Book Your Paradise.",
                 font_size=10, color=fg)


# ============================================================
# SLIDE 1: Title Slide (Teal background)
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
add_bg_rect(slide1, TEAL)

# Accent strip at top
add_bg_rect(slide1, CORAL, top=0, height=Inches(0.12))

# Logo area
add_logo_text(slide1, Inches(0.8), Inches(1.2), size=16)

# Title
add_text_box(slide1, Inches(0.8), Inches(2.5), Inches(10), Inches(1.5),
             "Presentation Title Here",
             font_size=44, color=WHITE, bold=True)

# Subtitle
add_text_box(slide1, Inches(0.8), Inches(4.2), Inches(8), Inches(0.8),
             "Subtitle or description goes here  |  Date",
             font_size=20, color=RGBColor(0xB0, 0xD8, 0xD2))

# Tagline at bottom
add_text_box(slide1, Inches(0.8), Inches(6.2), Inches(8), Inches(0.5),
             "Name Your Price. Book Your Paradise.",
             font_size=14, color=RGBColor(0x8C, 0xC5, 0xBC))

# Coral accent line
add_bg_rect(slide1, CORAL, left=Inches(0.8), top=Inches(5.8),
            width=Inches(3), height=Inches(0.05))

add_footer(slide1, dark_bg=True)


# ============================================================
# SLIDE 2: Section Divider (Teal header + Cream body)
# ============================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(slide2, CREAM)

# Teal header band
add_bg_rect(slide2, TEAL, height=Inches(2.8))
add_bg_rect(slide2, CORAL, top=Inches(2.8), height=Inches(0.08))

add_logo_text(slide2, Inches(0.8), Inches(0.5), size=12)

# Section title
add_text_box(slide2, Inches(0.8), Inches(1.2), Inches(10), Inches(1.2),
             "Section Title",
             font_size=40, color=WHITE, bold=True)

# Body text area placeholder
add_text_box(slide2, Inches(0.8), Inches(3.4), Inches(10), Inches(0.8),
             "Brief overview of what this section covers. Use this slide to introduce a new topic or group of related content.",
             font_size=18, color=NAVY)

add_footer(slide2)


# ============================================================
# SLIDE 3: Content Slide (Cream background, bullets)
# ============================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(slide3, CREAM)

# Top bar
add_bg_rect(slide3, TEAL, height=Inches(0.08))

# Title
add_text_box(slide3, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
             "Content Slide Title",
             font_size=32, color=TEAL, bold=True)

# Accent underline
add_accent_bar(slide3, Inches(1.25))

# Bullet points
bullets = [
    "First key point — keep it concise and actionable",
    "Second key point — one idea per bullet",
    "Third key point — use data to support claims",
    "Fourth key point — highlight with Coral for emphasis",
    "Fifth key point — end with a clear takeaway",
]

y_pos = Inches(1.7)
for i, bullet in enumerate(bullets):
    # Bullet dot
    dot = slide3.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(0.85), y_pos + Inches(0.12), Inches(0.15), Inches(0.15))
    dot.fill.solid()
    dot.fill.fore_color.rgb = CORAL if i == 3 else TEAL
    dot.line.fill.background()

    # Bullet text
    add_text_box(slide3, Inches(1.2), y_pos, Inches(10), Inches(0.5),
                 bullet, font_size=18, color=NAVY)
    y_pos += Inches(0.75)

add_footer(slide3)


# ============================================================
# SLIDE 4: Two-Column Layout (Stats / Features)
# ============================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(slide4, CREAM)
add_bg_rect(slide4, TEAL, height=Inches(0.08))

add_text_box(slide4, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
             "Key Metrics & Highlights",
             font_size=32, color=TEAL, bold=True)

add_accent_bar(slide4, Inches(1.25))

# Left column — stats
stats = [
    ("117+", "Resorts Worldwide"),
    ("351", "Unit Types Available"),
    ("34%", "Voice Search Adoption"),
    ("99.97%", "Platform Uptime"),
]

for i, (number, label) in enumerate(stats):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + col * Inches(5.5)
    y = Inches(1.8) + row * Inches(2.2)

    # Stat card background
    card = slide4.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(4.8), Inches(1.8))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.fill.background()

    # Coral top accent on card
    accent = slide4.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, x, y, Inches(4.8), Inches(0.06))
    accent.fill.solid()
    accent.fill.fore_color.rgb = CORAL
    accent.line.fill.background()

    # Number
    add_text_box(slide4, x + Inches(0.3), y + Inches(0.3), Inches(4), Inches(0.8),
                 number, font_size=42, color=TEAL, bold=True)

    # Label
    add_text_box(slide4, x + Inches(0.3), y + Inches(1.1), Inches(4), Inches(0.5),
                 label, font_size=16, color=MUTED)

add_footer(slide4)


# ============================================================
# SLIDE 5: Feature Highlight (Image placeholder + text)
# ============================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(slide5, WHITE)
add_bg_rect(slide5, TEAL, height=Inches(0.08))

# Left side — image placeholder
img_placeholder = slide5.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.8),
    Inches(5.5), Inches(5.8))
img_placeholder.fill.solid()
img_placeholder.fill.fore_color.rgb = SAND
img_placeholder.line.color.rgb = RGBColor(0xD0, 0xCB, 0xC3)
img_placeholder.line.width = Pt(1)

add_text_box(slide5, Inches(1.5), Inches(3.2), Inches(4), Inches(1),
             "[ Insert Image Here ]",
             font_size=18, color=MUTED, alignment=PP_ALIGN.CENTER)

# Right side — content
add_text_box(slide5, Inches(7), Inches(1.2), Inches(5.5), Inches(0.8),
             "Feature Highlight",
             font_size=32, color=TEAL, bold=True)

# Accent bar
shape = slide5.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(7), Inches(1.95), Inches(2), Inches(0.06))
shape.fill.solid()
shape.fill.fore_color.rgb = CORAL
shape.line.fill.background()

add_text_box(slide5, Inches(7), Inches(2.3), Inches(5.5), Inches(2),
             "Describe the feature or value proposition here. "
             "Keep it to 2-3 sentences that communicate the key benefit to your audience.",
             font_size=18, color=NAVY)

# Call-to-action box
cta_shape = slide5.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(4.8),
    Inches(3.5), Inches(0.7))
cta_shape.fill.solid()
cta_shape.fill.fore_color.rgb = CORAL
cta_shape.line.fill.background()

# CTA text
cta_tf = cta_shape.text_frame
cta_tf.paragraphs[0].text = "Call to Action"
cta_tf.paragraphs[0].font.size = Pt(18)
cta_tf.paragraphs[0].font.color.rgb = WHITE
cta_tf.paragraphs[0].font.bold = True
cta_tf.paragraphs[0].font.name = "Roboto"
cta_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
cta_tf.word_wrap = True

add_footer(slide5)


# ============================================================
# SLIDE 6: Closing / Thank You (Teal background)
# ============================================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(slide6, TEAL)
add_bg_rect(slide6, CORAL, top=0, height=Inches(0.12))

add_logo_text(slide6, Inches(0.8), Inches(1.0), size=16)

# Thank you
add_text_box(slide6, Inches(0.8), Inches(2.5), Inches(10), Inches(1.2),
             "Thank You",
             font_size=52, color=WHITE, bold=True)

# Contact info
add_text_box(slide6, Inches(0.8), Inches(4.0), Inches(8), Inches(0.5),
             "rent-a-vacation.com",
             font_size=22, color=RGBColor(0xB0, 0xD8, 0xD2))

add_text_box(slide6, Inches(0.8), Inches(4.6), Inches(8), Inches(0.5),
             "support@rent-a-vacation.com  |  1-800-RAV-BOOK",
             font_size=16, color=RGBColor(0x8C, 0xC5, 0xBC))

# Tagline
add_bg_rect(slide6, CORAL, left=Inches(0.8), top=Inches(5.6),
            width=Inches(3), height=Inches(0.05))

add_text_box(slide6, Inches(0.8), Inches(5.8), Inches(8), Inches(0.5),
             "Name Your Price. Book Your Paradise.",
             font_size=16, color=RGBColor(0x8C, 0xC5, 0xBC))

add_footer(slide6, dark_bg=True)


# ============================================================
# Save
# ============================================================
output_path = "docs/RAV-Brand-Template.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
print("Slide overview:")
print("  1. Title Slide (teal bg)")
print("  2. Section Divider (teal header + cream body)")
print("  3. Content Slide (bullets)")
print("  4. Key Metrics (2x2 stat cards)")
print("  5. Feature Highlight (image + text split)")
print("  6. Closing / Thank You (teal bg)")
